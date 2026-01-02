import streamlit as st
import fitz  # PyMuPDF
import numpy as np
from sentence_transformers import SentenceTransformer
from PIL import Image, ImageDraw
from io import BytesIO
from pathlib import Path
import docx
import pptx
import difflib
import cv2

#  CONFIG 
SIM_THRESHOLD = 0.80
MAX_DIFF_LINES = 3
IMAGE_SIM_THRESHOLD = 10  # lower means more similar

#  MODEL 
@st.cache_resource
def load_model():
    return SentenceTransformer("all-MiniLM-L6-v2")

model = load_model()

#  UTILS 
def clean_text(text):
    if not text:
        return ""
    return " ".join(text.replace("\n", " ").replace("\r", " ").split())

def cosine_similarity(a, b):
    return float(np.dot(a, b))

def extract_text_diff(a, b):
    a_lines = a.split(". ")
    b_lines = b.split(". ")
    diff = list(difflib.ndiff(a_lines, b_lines))
    removed = [l[2:] for l in diff if l.startswith("- ")][:MAX_DIFF_LINES]
    added = [l[2:] for l in diff if l.startswith("+ ")][:MAX_DIFF_LINES]
    return (
        " | ".join(removed) if removed else "No clear removed text",
        " | ".join(added) if added else "No clear added text"
    )

# IMAGE COMPARISON
def compare_page_images(page_a_img, page_b_img):
    img1 = np.array(page_a_img.convert("L"))
    img2 = np.array(page_b_img.convert("L"))

    if img1.shape != img2.shape:
        img2 = cv2.resize(img2, (img1.shape[1], img1.shape[0]))

    diff = cv2.absdiff(img1, img2)
    _, thresh = cv2.threshold(diff, 30, 255, cv2.THRESH_BINARY)

    kernel = np.ones((5,5), np.uint8)
    thresh = cv2.dilate(thresh, kernel, iterations=2)

    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    highlighted = page_b_img.copy()
    draw = ImageDraw.Draw(highlighted)

    for cnt in contours:
        x, y, w, h = cv2.boundingRect(cnt)
        if w * h > 500:
            draw.rectangle([x, y, x+w, y+h], outline="red", width=3)

    return highlighted

def extract_pdf_images(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    images = []
    for page in doc:
        pix = page.get_pixmap(dpi=150)
        img = Image.open(BytesIO(pix.tobytes("png"))).convert("RGB")
        images.append(img)
    return images

#  EXTRACTION
def extract_pdf(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        images_count = len(page.get_images(full=True))

        pages.append({
            "page": page_num,
            "text": clean_text(text),   # empty for scanned PDFs
            "images": images_count
        })

    return pages, doc.page_count

def extract_docx(file):
    file.seek(0)
    doc = docx.Document(file)
    text = clean_text(" ".join(p.text for p in doc.paragraphs))
    return [{"page": 1, "text": text, "images": 0}], 1

def extract_txt(file):
    file.seek(0)
    raw = file.read()
    try:
        text = raw.decode("utf-8")
    except:
        text = raw.decode("latin-1")
    return [{"page": 1, "text": clean_text(text), "images": 0}], 1


def extract_ppt(file):
    file.seek(0)
    prs = pptx.Presentation(BytesIO(file.read()))

    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        image_count = 0

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if shape.shape_type == 13:
                image_count += 1

        slides.append({
            "page": i,
            "text": clean_text(" ".join(slide_text)),
            "images": image_count
        })

    return slides, len(slides)

def extract_image(file):
    file.seek(0)
    return Image.open(file).convert("RGB")

def image_similarity(img1, img2):
    img1 = cv2.resize(np.array(img1), (256, 256))
    img2 = cv2.resize(np.array(img2), (256, 256))
    diff = np.mean(np.abs(img1 - img2))
    return diff < IMAGE_SIM_THRESHOLD

def image_diff(img1, img2):
    img1 = cv2.resize(np.array(img1), (512, 512))
    img2 = cv2.resize(np.array(img2), (512, 512))
    diff = cv2.absdiff(img1, img2)
    return Image.fromarray(diff)

def extract(upload):
    ext = Path(upload.name).suffix.lower()

    if ext == ".pdf":
        return extract_pdf(upload.read())

    if ext == ".docx":
        return extract_docx(upload)

    if ext == ".txt":
        return extract_txt(upload)

    if ext == ".ppt":
        st.error("❌ .ppt (old PowerPoint) is not supported. Please upload .pptx")
        st.stop()

    if ext == ".pptx":
        return extract_ppt(upload)

    if ext in [".jpg", ".jpeg", ".png"]:
        return extract_image(upload), "image"

    return [], 0

#  COMPARISON 
def compare_pages(pages_a, pages_b):
    changes = []
    similarities = []

    max_pages = max(len(pages_a), len(pages_b))

    for i in range(max_pages):
        page_no = i + 1

        if i >= len(pages_a):
            changes.append({
                "Page": page_no,
                "Change Type": "Page Added",
                "Before": "-",
                "After": "New page added"
            })
            continue

        if i >= len(pages_b):
            changes.append({
                "Page": page_no,
                "Change Type": "Page Removed",
                "Before": "Page existed",
                "After": "-"
            })
            continue

        a = pages_a[i]
        b = pages_b[i]

        # Visual count comparison (CRITICAL for scanned PDFs)
        if b["images"] > a["images"]:
            changes.append({
                "Page": page_no,
                "Change Type": "Visual Added",
                "Before": f"{a['images']} image(s)",
                "After": f"{b['images']} image(s)"
            })
        elif b["images"] < a["images"]:
            changes.append({
                "Page": page_no,
                "Change Type": "Visual Removed",
                "Before": f"{a['images']} image(s)",
                "After": f"{b['images']} image(s)"
            })

        # Text similarity (only if text exists)
        if a["text"] and b["text"]:
            emb = model.encode([a["text"], b["text"]], normalize_embeddings=True)
            sim = cosine_similarity(emb[0], emb[1])
            similarities.append(sim)

            if sim < SIM_THRESHOLD:
                removed, added = extract_text_diff(a["text"], b["text"])
                changes.append({
                    "Page": page_no,
                    "Change Type": "Content Modified",
                    "Before": removed,
                    "After": added
                })

    avg_similarity = sum(similarities) / max(len(similarities), 1)
    return changes, avg_similarity

#  STREAMLIT UI 
st.set_page_config("Document Comparison", layout="wide")
st.title("📄 Document Comparison")

col1, col2 = st.columns(2)
with col1:
    file_a = st.file_uploader(
        "Upload Document A (Original)",
        type=["pdf","docx","txt","pptx","jpg","png"]
    )
with col2:
    file_b = st.file_uploader(
        "Upload Document B (Updated)",
        type=["pdf","docx","txt","pptx","jpg","png"]
    )

if st.button("Compare Documents"):
    if not file_a or not file_b:
        st.error("Please upload both documents.")
    else:
        # IMAGE-ONLY COMPARISON
        if file_a.name.lower().endswith((".jpg", ".jpeg", ".png")) and \
           file_b.name.lower().endswith((".jpg", ".jpeg", ".png")):

            img_a = extract_image(file_a)
            img_b = extract_image(file_b)

            images_similar = "YES" if image_similarity(img_a, img_b) else "NO"
            st.metric("Images Similar?", images_similar)

            c1, c2, c3 = st.columns(3)
            c1.image(img_a, caption="Original Image", width=400)
            c2.image(img_b, caption="Updated Image", width=400)

            if images_similar == "NO":
                c3.image(image_diff(img_a, img_b), caption="Difference", width=400)
            else:
                c3.success("No visible difference")

            st.stop()

        with st.spinner("Comparing documents..."):
            pages_a, count_a = extract(file_a)
            pages_b, count_b = extract(file_b)
            detected_changes, similarity = compare_pages(pages_a, pages_b)

            pages_a_imgs = []
            pages_b_imgs = []

            if file_a.name.lower().endswith(".pdf") and file_b.name.lower().endswith(".pdf"):
                file_a.seek(0)
                pages_a_imgs = extract_pdf_images(file_a.read())
                file_b.seek(0)
                pages_b_imgs = extract_pdf_images(file_b.read())

        documents_similar = "YES" if not detected_changes and count_a == count_b else "NO"

        st.subheader("Result")
        st.metric("Documents Similar?", documents_similar)

        st.subheader("Page Count")
        st.write(f"Document A: {count_a} pages")
        st.write(f"Document B: {count_b} pages")

        st.subheader("Detected Changes (Side-by-Side)")
        if detected_changes:
            for change in detected_changes:
                st.markdown(f"**Page {change['Page']} - {change['Change Type']}**")
                c1, c2 = st.columns(2)
                c1.markdown("**Original / Before:**")
                c1.write(change["Before"])
                c2.markdown("**Updated / After:**")
                c2.write(change["After"])
                st.markdown("---")
        else:
            st.success("No significant changes detected.")

        st.subheader("Visual Differences (Added/Removed Signatures or Images)")
        for i in range(min(len(pages_a_imgs), len(pages_b_imgs))):
            if pages_b[i]["images"] != pages_a[i]["images"]:
                highlighted = compare_page_images(pages_a_imgs[i], pages_b_imgs[i])
                st.image(highlighted, caption=f"Page {i+1} visual differences", width=600)
st.success("No significant changes detected.")

